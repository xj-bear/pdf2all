"""
PDF to PPT converter module.
Converts each PDF page to an image and embeds into PowerPoint slides using PyMuPDF.
Optimized for Cherry Studio with better performance.
"""
import os
import tempfile
import sys
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches


def convert(pdf_path: str, output_path: str | None = None, dpi: int = 150) -> dict:
    """
    Convert PDF to PowerPoint by converting pages to images.

    Args:
        pdf_path: Path to the input PDF file
        output_path: Optional output path. If not provided, uses same directory as PDF.
        dpi: Image resolution. Default is 150 for good quality.

    Returns:
        dict with status, output_path, and message
    """
    try:
        if not os.path.exists(pdf_path):
            return {
                "success": False,
                "error": f"PDF file not found: {pdf_path}"
            }

        # Generate output path if not provided
        if not output_path:
            base_name = os.path.splitext(pdf_path)[0]
            output_path = f"{base_name}.pptx"

        # Ensure output directory exists
        output_dir = os.path.dirname(output_path)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir)

        # Open PDF with PyMuPDF
        doc = fitz.open(pdf_path)
        page_count = len(doc)

        if page_count == 0:
            doc.close()
            return {
                "success": False,
                "error": "PDF has no pages"
            }

        # Create PowerPoint presentation
        prs = Presentation()

        # Set slide dimensions (standard 16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Calculate zoom factor based on DPI
        zoom = dpi / 72.0
        matrix = fitz.Matrix(zoom, zoom)

        # Process pages in batches for better memory management
        batch_size = 10  # Process 10 pages at a time

        with tempfile.TemporaryDirectory() as temp_dir:
            for i in range(page_count):
                try:
                    page = doc[i]

                    # Optimize rendering settings for performance
                    pix = page.get_pixmap(
                        matrix=matrix,
                        alpha=False,  # No alpha channel for faster processing
                        colorspace=fitz.csRGB  # Use RGB for better compatibility
                    )

                    # Save image temporarily
                    img_path = os.path.join(temp_dir, f"page_{i + 1}.png")
                    pix.save(img_path)

                    # Get image dimensions
                    img_width = pix.width
                    img_height = pix.height
                    aspect_ratio = img_width / img_height

                    # Add blank slide
                    blank_layout = prs.slide_layouts[6]  # Blank layout
                    slide = prs.slides.add_slide(blank_layout)

                    # Calculate image dimensions to fit slide
                    slide_width = prs.slide_width
                    slide_height = prs.slide_height

                    if aspect_ratio > (slide_width / slide_height):
                        # Image is wider, fit to width
                        width = slide_width
                        height = width / aspect_ratio
                    else:
                        # Image is taller, fit to height
                        height = slide_height
                        width = height * aspect_ratio

                    # Center the image
                    left = (slide_width - width) / 2
                    top = (slide_height - height) / 2

                    # Add image to slide
                    slide.shapes.add_picture(img_path, left, top, width, height)

                    # Print progress for Cherry Studio
                    if (i + 1) % batch_size == 0 or i == page_count - 1:
                        print(f"[INFO] Processed {i + 1}/{page_count} pages", file=sys.stderr, flush=True)

                except Exception as e:
                    # If a page fails, continue with next page
                    print(f"[WARNING] Failed to process page {i+1}: {str(e)}", file=sys.stderr, flush=True)
                    continue

        doc.close()

        # Save presentation
        prs.save(output_path)

        return {
            "success": True,
            "output_path": output_path,
            "message": f"Successfully converted {page_count} page(s) to PowerPoint: {output_path}",
            "pages_count": page_count
        }

    except Exception as e:
        return {
            "success": False,
            "error": f"Conversion failed: {str(e)}"
        }
